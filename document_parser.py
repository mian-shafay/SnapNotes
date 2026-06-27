import pymupdf4llm
from pptx import Presentation
import tempfile
import os

def parse_pdf(file_bytes: bytes) -> str:
    """Extracts text from PDF and converts it to Markdown, preserving headings."""
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name

    try:
        # pymupdf4llm is great for preserving document structure and headings as markdown
        md_text = pymupdf4llm.to_markdown(tmp_path)
        return md_text
    finally:
        os.remove(tmp_path)

def parse_pptx(file_bytes: bytes) -> str:
    """Extracts text from PPTX, using Slide titles as Headings."""
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
        tmp.write(file_bytes)
        tmp_path = tmp.name
        
    try:
        prs = Presentation(tmp_path)
        md_text = ""
        
        for i, slide in enumerate(prs.slides):
            title = ""
            if slide.shapes.title and slide.shapes.title.text:
                title = slide.shapes.title.text.strip()
                md_text += f"\n## {title}\n\n"
            else:
                md_text += f"\n## Slide {i+1}\n\n"
            
            for shape in slide.shapes:
                # Skip the title shape as we already added it
                if shape == slide.shapes.title:
                    continue
                if hasattr(shape, "text") and shape.text.strip():
                    md_text += f"{shape.text.strip()}\n"
            
            md_text += "\n---\n" # Slide separator
            
        return md_text
    finally:
        os.remove(tmp_path)
